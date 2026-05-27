#!/usr/bin/env python3
"""
CueDeck demo presentation generator.
Creates demo.pptx → converts to demo.pdf via LibreOffice.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import subprocess, os, sys

# ── Colors ──────────────────────────────────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BG      = RGBColor(0xF9, 0xFA, 0xFB)   # almost-white background
DARK    = RGBColor(0x11, 0x18, 0x27)   # near-black
GRAY    = RGBColor(0x6B, 0x72, 0x80)   # subtitle / caption
BLUE    = RGBColor(0x25, 0x63, 0xEB)   # accent
BLUE_LT = RGBColor(0xDB, 0xEA, 0xFE)  # light-blue chip bg
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
RED     = RGBColor(0xDC, 0x26, 0x26)
AMBER   = RGBColor(0xD9, 0x77, 0x06)
LINE    = RGBColor(0xE5, 0xE7, 0xEB)   # divider

# ── Helpers ──────────────────────────────────────────────────────────────────
W = Inches(13.33)   # 16:9 widescreen width
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_color, radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=24, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "SF Pro Display"
    return tb

def hline(slide, y, color=LINE):
    line = slide.shapes.add_connector(1, Inches(0.5), y, W - Inches(0.5), y)
    line.line.color.rgb = color
    line.line.width = Pt(0.75)

def chip(slide, text, x, y, size=16,
         bg_color=BLUE_LT, text_color=BLUE):
    w = Inches(len(text) * 0.11 + 0.5)
    h = Inches(0.38)
    r = rect(slide, x, y, w, h, bg_color)
    txt(slide, text, x, y + Inches(0.04), w, h,
        size=size, color=text_color, align=PP_ALIGN.CENTER)
    return w

# ── Slides ───────────────────────────────────────────────────────────────────

def slide_title(prs):
    """01 — Title"""
    sl = blank(prs)
    bg(sl, WHITE)

    # left accent bar
    rect(sl, Inches(0), Inches(0), Inches(0.08), H, BLUE)

    # big title
    txt(sl, "CueDeck", Inches(0.6), Inches(1.8), Inches(8), Inches(1.5),
        size=80, bold=True, color=DARK)

    # tagline
    txt(sl, "Управление презентациями для живых мероприятий",
        Inches(0.6), Inches(3.4), Inches(9), Inches(0.7),
        size=28, color=GRAY)

    # version chip
    chip(sl, "v 0.1.2", Inches(0.6), Inches(4.4), size=18)

    # right decorative block
    rect(sl, Inches(10.2), Inches(0), Inches(3.13), H, BLUE)
    txt(sl, "🎬", Inches(10.8), Inches(2.8), Inches(2), Inches(2),
        size=80, color=WHITE, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    """02 — Problem"""
    sl = blank(prs)
    bg(sl)

    txt(sl, "Проблема", Inches(0.7), Inches(0.4), Inches(11), Inches(0.6),
        size=14, color=GRAY)
    txt(sl, "Один оператор — три экрана", Inches(0.7), Inches(0.9), Inches(11), Inches(1.0),
        size=44, bold=True, color=DARK)
    hline(sl, Inches(2.1))

    items = [
        ("📋", "Плейлист из PDF нескольких спикеров"),
        ("🖥", "Зрительный зал — чистый слайд, без лишнего"),
        ("🪞", "Спикер видит заметки и таймер, не отвлекаясь"),
        ("⏱", "Оператор контролирует время прямо с пульта"),
    ]
    for i, (icon, line) in enumerate(items):
        y = Inches(2.4 + i * 1.1)
        txt(sl, icon, Inches(0.7), y, Inches(0.6), Inches(0.9), size=28)
        txt(sl, line, Inches(1.5), y + Inches(0.1), Inches(10), Inches(0.8), size=24, color=DARK)


def slide_three_roles(prs):
    """03 — Three roles"""
    sl = blank(prs)
    bg(sl)

    txt(sl, "Три роли — три окна", Inches(0.7), Inches(0.35), Inches(11), Inches(0.8),
        size=44, bold=True, color=DARK)
    hline(sl, Inches(1.35))

    roles = [
        (BLUE,               "🎛",  "Оператор",  "Управляет слайдами\nи видит всё сразу"),
        (RGBColor(0x7C,0x3A,0xED), "🎤", "Спикер",   "Суфлёр: следующий\nслайд + таймер + notes"),
        (RGBColor(0x0F,0x76,0x6E), "👥", "Аудитория","Чистый fullscreen\nна внешнем дисплее"),
    ]
    col_w = Inches(4.0)
    for i, (color, icon, title, desc) in enumerate(roles):
        x = Inches(0.5 + i * 4.27)
        y = Inches(1.65)
        rect(sl, x, y, col_w, Inches(4.8), color)
        txt(sl, icon, x, y + Inches(0.4), col_w, Inches(1.0),
            size=44, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, title, x, y + Inches(1.4), col_w, Inches(0.7),
            size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, desc, x, y + Inches(2.15), col_w, Inches(1.4),
            size=18, color=WHITE, align=PP_ALIGN.CENTER)


def slide_operator(prs):
    """04 — Operator view"""
    sl = blank(prs)
    bg(sl)

    txt(sl, "Operator View", Inches(0.7), Inches(0.35), Inches(11), Inches(0.8),
        size=44, bold=True, color=DARK)
    hline(sl, Inches(1.35))

    # mock operator panel
    panel_x, panel_y = Inches(0.6), Inches(1.6)
    panel_w, panel_h = Inches(12.1), Inches(5.4)
    rect(sl, panel_x, panel_y, panel_w, panel_h, RGBColor(0x1E,0x29,0x3B))

    # current slide area
    rect(sl, panel_x + Inches(0.2), panel_y + Inches(0.2),
         Inches(5.5), Inches(3.1), RGBColor(0x33,0x41,0x55))
    txt(sl, "Текущий слайд", panel_x + Inches(0.2),
        panel_y + Inches(0.2), Inches(5.5), Inches(3.1),
        size=20, color=RGBColor(0xE2,0xE8,0xF0), align=PP_ALIGN.CENTER)

    # next slide
    rect(sl, panel_x + Inches(5.9), panel_y + Inches(0.2),
         Inches(2.6), Inches(1.5), RGBColor(0x33,0x41,0x55))
    txt(sl, "Следующий", panel_x + Inches(5.9),
        panel_y + Inches(0.2), Inches(2.6), Inches(1.5),
        size=14, color=RGBColor(0xE2,0xE8,0xF0), align=PP_ALIGN.CENTER)

    # timer
    rect(sl, panel_x + Inches(8.7), panel_y + Inches(0.2),
         Inches(3.2), Inches(1.5), RGBColor(0x16,0xA3,0x4A))
    txt(sl, "12:34", panel_x + Inches(8.7), panel_y + Inches(0.35),
        Inches(3.2), Inches(1.0),
        size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # playlist items
    playlist = ["📄 Иванов — Введение.pdf", "📄 Петров — Кейсы.pdf", "📄 Сидорова — Итоги.pdf"]
    for i, name in enumerate(playlist):
        item_y = panel_y + Inches(0.2) + Inches(i * 0.6) + Inches(3.4)
        bg_c = RGBColor(0x3B,0x82,0xF6) if i == 0 else RGBColor(0x2D,0x3D,0x52)
        rect(sl, panel_x + Inches(0.2), item_y, Inches(11.7), Inches(0.5), bg_c)
        txt(sl, name, panel_x + Inches(0.4), item_y + Inches(0.05),
            Inches(11), Inches(0.45), size=14,
            color=WHITE if i == 0 else RGBColor(0xCB,0xD5,0xE1))


def slide_speaker(prs):
    """05 — Speaker view"""
    sl = blank(prs)
    bg(sl)

    txt(sl, "Speaker View  —  суфлёр для спикера", Inches(0.7), Inches(0.35),
        Inches(11), Inches(0.8), size=44, bold=True, color=DARK)
    hline(sl, Inches(1.35))

    # mock speaker window (dark)
    mock_x, mock_y = Inches(0.6), Inches(1.65)
    mock_w, mock_h = Inches(12.1), Inches(5.3)
    rect(sl, mock_x, mock_y, mock_w, mock_h, RGBColor(0x0D,0x11,0x17))

    # current slide preview
    rect(sl, mock_x + Inches(0.3), mock_y + Inches(0.3),
         Inches(5.2), Inches(2.9), RGBColor(0x1C,0x2B,0x3A))
    txt(sl, "Текущий слайд", mock_x + Inches(0.3), mock_y + Inches(0.3),
        Inches(5.2), Inches(2.9), size=18,
        color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

    # next slide
    rect(sl, mock_x + Inches(5.7), mock_y + Inches(0.3),
         Inches(2.5), Inches(2.9), RGBColor(0x1C,0x2B,0x3A))
    txt(sl, "Следующий", mock_x + Inches(5.7), mock_y + Inches(0.3),
        Inches(2.5), Inches(2.9), size=14,
        color=RGBColor(0x6B,0x72,0x80), align=PP_ALIGN.CENTER)

    # timer overlay (top-right corner)
    rect(sl, mock_x + Inches(9.2), mock_y + Inches(0.2),
         Inches(2.7), Inches(0.85), RGBColor(0x16,0xA3,0x4A))
    txt(sl, "12:34", mock_x + Inches(9.2), mock_y + Inches(0.2),
        Inches(2.7), Inches(0.85), size=36, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

    # notes area
    rect(sl, mock_x + Inches(0.3), mock_y + Inches(3.35),
         Inches(11.5), Inches(1.7), RGBColor(0x1C,0x2B,0x3A))
    txt(sl, "📝  Здесь видны заметки к текущему слайду — только у спикера",
        mock_x + Inches(0.5), mock_y + Inches(3.5),
        Inches(11.0), Inches(1.3), size=16,
        color=RGBColor(0xF1,0xF5,0xF9))


def slide_timer(prs):
    """06 — Timer modes"""
    sl = blank(prs)
    bg(sl)

    txt(sl, "Таймер", Inches(0.7), Inches(0.35), Inches(11), Inches(0.8),
        size=44, bold=True, color=DARK)
    hline(sl, Inches(1.35))

    modes = [
        (GREEN,  "⏳",  "Countdown",  "Обратный отсчёт\nОкраска: зелёный → жёлтый → красный"),
        (BLUE,   "⏱",  "Stopwatch",  "Счётчик вверх\nОт нуля"),
        (GRAY,   "🕐",  "Clock",      "Системное время\nНейтральный цвет"),
    ]
    for i, (color, icon, title, desc) in enumerate(modes):
        x = Inches(0.5 + i * 4.27)
        y = Inches(1.65)
        w = Inches(4.0)
        rect(sl, x, y, w, Inches(4.7), color)
        txt(sl, icon, x, y + Inches(0.35), w, Inches(1.0),
            size=48, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, title, x, y + Inches(1.4), w, Inches(0.7),
            size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(sl, desc, x, y + Inches(2.15), w, Inches(1.8),
            size=16, color=WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "Позиция: угол экрана  ·  Масштаб: 0.5× — 2.5×  ·  Только на экране спикера",
        Inches(0.7), Inches(6.85), Inches(12), Inches(0.5),
        size=14, color=GRAY)


def slide_blackout(prs):
    """07 — Blackout & Key Visual"""
    sl = blank(prs)
    bg(sl)

    txt(sl, "Blackout & Key Visual", Inches(0.7), Inches(0.35), Inches(11), Inches(0.8),
        size=44, bold=True, color=DARK)
    hline(sl, Inches(1.35))

    # blackout mock
    bx, by = Inches(0.6), Inches(1.65)
    rect(sl, bx, by, Inches(5.8), Inches(4.5), DARK)
    txt(sl, "⬛  Blackout", bx, by + Inches(1.8), Inches(5.8), Inches(1.0),
        size=28, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    txt(sl, "Экран зала — чёрный\nСлайды у оператора — без изменений",
        bx + Inches(0.3), by + Inches(2.9), Inches(5.2), Inches(1.2),
        size=15, color=GRAY, align=PP_ALIGN.CENTER)

    # key visual mock
    kx = Inches(7.0)
    rect(sl, kx, by, Inches(5.8), Inches(4.5), RGBColor(0x1E,0x29,0x3B))
    # simulate logo placeholder
    rect(sl, kx + Inches(1.5), by + Inches(0.9), Inches(2.8), Inches(1.6),
         RGBColor(0x3B,0x82,0xF6))
    txt(sl, "ЛОГОТИП\nКОМПАНИИ", kx + Inches(1.5), by + Inches(0.9),
        Inches(2.8), Inches(1.6),
        size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "Key Visual: пока идёт перерыв\nзрители видят брендинг",
        kx + Inches(0.3), by + Inches(2.9), Inches(5.2), Inches(1.2),
        size=15, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

    # labels
    txt(sl, "Нажать B", Inches(2.7), Inches(6.3), Inches(2), Inches(0.5),
        size=14, color=GRAY, align=PP_ALIGN.CENTER)
    txt(sl, "Blackout + Key Visual", Inches(8.9), Inches(6.3), Inches(3), Inches(0.5),
        size=14, color=GRAY, align=PP_ALIGN.CENTER)


def slide_playlist(prs):
    """08 — Playlist"""
    sl = blank(prs)
    bg(sl)

    txt(sl, "Плейлист", Inches(0.7), Inches(0.35), Inches(11), Inches(0.8),
        size=44, bold=True, color=DARK)
    hline(sl, Inches(1.35))

    txt(sl, "Несколько спикеров — один проект",
        Inches(0.7), Inches(1.5), Inches(11), Inches(0.6),
        size=24, color=GRAY)

    entries = [
        ("01", "Иванова А.В.",   "Введение_2026.pdf",        "28 слайдов", BLUE),
        ("02", "Петров К.С.",    "Кейс_Сбер.pdf",            "14 слайдов", RGBColor(0x7C,0x3A,0xED)),
        ("03", "Сидорова М.",    "Итоги_квартал.pdf",        "22 слайда",  RGBColor(0x0F,0x76,0x6E)),
        ("04", "Хусаенов А.",    "Закрытие_2026.pptx",       "10 слайдов", RGBColor(0xD9,0x77,0x06)),
    ]
    for i, (num, speaker, fname, slides, color) in enumerate(entries):
        y = Inches(2.3 + i * 1.1)
        rect(sl, Inches(0.6), y, Inches(0.5), Inches(0.85), color)
        txt(sl, num, Inches(0.6), y + Inches(0.15), Inches(0.5), Inches(0.6),
            size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        rect(sl, Inches(1.2), y, Inches(11.5), Inches(0.85),
             RGBColor(0xF3,0xF4,0xF6) if i % 2 == 0 else WHITE)
        txt(sl, speaker, Inches(1.4), y + Inches(0.15), Inches(3.5), Inches(0.6),
            size=18, bold=True, color=DARK)
        txt(sl, fname, Inches(4.8), y + Inches(0.15), Inches(5.0), Inches(0.6),
            size=16, color=GRAY)
        txt(sl, slides, Inches(11.0), y + Inches(0.15), Inches(1.5), Inches(0.6),
            size=14, color=BLUE)


def slide_formats(prs):
    """09 — Supported formats"""
    sl = blank(prs)
    bg(sl)

    txt(sl, "Форматы файлов", Inches(0.7), Inches(0.35), Inches(11), Inches(0.8),
        size=44, bold=True, color=DARK)
    hline(sl, Inches(1.35))

    formats = [
        ("📄", "PDF",       "Нативный",    GREEN,  "Рендер в каждом окне независимо"),
        ("📊", "PPTX / PPT","via LibreOffice", BLUE, "Кешируется по SHA1"),
        ("📑", "ODP / KEY", "via LibreOffice", BLUE, "OpenDocument + Keynote"),
        ("🖼", "PNG / JPG", "Нативный",    GREEN,  "WEBP, GIF, BMP — тоже"),
    ]
    for i, (icon, fmt, method, m_color, note) in enumerate(formats):
        y = Inches(1.75 + i * 1.25)

        txt(sl, icon, Inches(0.6), y, Inches(0.7), Inches(1.0), size=32)
        txt(sl, fmt, Inches(1.4), y + Inches(0.15), Inches(2.5), Inches(0.7),
            size=24, bold=True, color=DARK)
        chip(sl, method, Inches(3.9), y + Inches(0.2), size=14,
             bg_color=BLUE_LT if m_color == BLUE else RGBColor(0xDC,0xFC,0xE7),
             text_color=m_color)
        txt(sl, note, Inches(6.3), y + Inches(0.15), Inches(6.5), Inches(0.7),
            size=16, color=GRAY)

        if i < 3:
            hline(sl, y + Inches(1.1), LINE)


def slide_notes(prs):
    """10 — Notes"""
    sl = blank(prs)
    bg(sl)

    txt(sl, "Заметки к слайдам", Inches(0.7), Inches(0.35), Inches(11), Inches(0.8),
        size=44, bold=True, color=DARK)
    hline(sl, Inches(1.35))

    items = [
        ("✏️", "Редактируются прямо в Operator View"),
        ("💾", "Сохраняются в .notes.json рядом с PDF — файл не трогается"),
        ("🔒", "Привязаны по SHA1: переименуй файл — заметки не потеряются"),
        ("👁", "Видны только на экране спикера, зрители не видят"),
        ("🔠", "Размер шрифта в суфлёре регулируется отдельно"),
    ]
    for i, (icon, line) in enumerate(items):
        y = Inches(1.75 + i * 1.0)
        txt(sl, icon, Inches(0.7), y, Inches(0.6), Inches(0.9), size=26)
        txt(sl, line, Inches(1.5), y + Inches(0.1), Inches(11), Inches(0.75),
            size=22, color=DARK)


def slide_install(prs):
    """11 — Installation"""
    sl = blank(prs)
    bg(sl)

    txt(sl, "Установка", Inches(0.7), Inches(0.35), Inches(11), Inches(0.8),
        size=44, bold=True, color=DARK)
    hline(sl, Inches(1.35))

    # macOS
    rect(sl, Inches(0.6), Inches(1.65), Inches(5.9), Inches(5.3), RGBColor(0xF0,0xF9,0xFF))
    txt(sl, "🍎  macOS", Inches(0.85), Inches(1.9), Inches(5.4), Inches(0.7),
        size=26, bold=True, color=DARK)
    mac_steps = [
        "Скачать .dmg с GitHub Releases",
        "Перетащить CueDeck в Applications",
        "При первом запуске:\nSystem Settings → Privacy → Open Anyway",
        "Готово — без регистрации, без подписки",
    ]
    for i, step in enumerate(mac_steps):
        txt(sl, f"  {i+1}.  {step}",
            Inches(0.85), Inches(2.7 + i * 0.9), Inches(5.4), Inches(0.85),
            size=16, color=DARK)

    # requirements
    rect(sl, Inches(7.0), Inches(1.65), Inches(5.9), Inches(5.3),
         RGBColor(0xFD,0xF4,0xFF))
    txt(sl, "📦  Опционально", Inches(7.25), Inches(1.9), Inches(5.4), Inches(0.7),
        size=26, bold=True, color=DARK)
    txt(sl, "LibreOffice\n— для открытия PPTX / ODP / KEY\n\nbrew install --cask libreoffice",
        Inches(7.25), Inches(2.65), Inches(5.4), Inches(2.0),
        size=16, color=DARK)
    txt(sl, "PDF и изображения работают\nбез LibreOffice",
        Inches(7.25), Inches(4.9), Inches(5.4), Inches(0.9),
        size=14, color=GRAY)


def slide_final(prs):
    """12 — Final"""
    sl = blank(prs)
    bg(sl, WHITE)

    rect(sl, Inches(0), Inches(0), Inches(0.08), H, BLUE)

    txt(sl, "CueDeck", Inches(0.6), Inches(1.8), Inches(8), Inches(1.2),
        size=72, bold=True, color=DARK)
    txt(sl, "Open Source · macOS · Бесплатно",
        Inches(0.6), Inches(3.1), Inches(9), Inches(0.6),
        size=24, color=GRAY)

    txt(sl, "github.com/ttpa3dhuk/PDF-Presenter",
        Inches(0.6), Inches(4.1), Inches(10), Inches(0.6),
        size=22, bold=True, color=BLUE)

    chip(sl, "v 0.1.2", Inches(0.6), Inches(5.0), size=18)

    rect(sl, Inches(10.2), Inches(0), Inches(3.13), H, BLUE)
    txt(sl, "✨", Inches(10.8), Inches(2.8), Inches(2), Inches(2),
        size=80, color=WHITE, align=PP_ALIGN.CENTER)


# ── Build ─────────────────────────────────────────────────────────────────────
def build():
    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_three_roles(prs)
    slide_operator(prs)
    slide_speaker(prs)
    slide_timer(prs)
    slide_blackout(prs)
    slide_playlist(prs)
    slide_formats(prs)
    slide_notes(prs)
    slide_install(prs)
    slide_final(prs)

    out_pptx = os.path.join(os.path.dirname(__file__), "cue-deck-demo.pptx")
    prs.save(out_pptx)
    print(f"✅  PPTX → {out_pptx}")

    # Convert to PDF via LibreOffice
    out_dir = os.path.dirname(out_pptx)
    result = subprocess.run(
        ["/opt/homebrew/bin/soffice", "--headless", "--convert-to", "pdf",
         "--outdir", out_dir, out_pptx],
        capture_output=True, text=True
    )
    if result.returncode == 0:
        pdf_path = out_pptx.replace(".pptx", ".pdf")
        print(f"✅  PDF  → {pdf_path}")
    else:
        print("❌  LibreOffice error:", result.stderr)
        sys.exit(1)

if __name__ == "__main__":
    build()
